"""
Open Manus Presentation Tool
Replicates the Manus.im `slides` tool functionality.
Generates .pptx presentations from a Markdown content outline.
"""
import json
import os
import re
import tempfile
from tools.registry import registry

def check_requirements() -> bool:
    try:
        import pptx
        return True
    except ImportError:
        return False

def create_presentation(content_markdown: str, output_filename: str = "presentation.pptx",
                        theme: str = "default", task_id: str = None) -> str:
    """
    Generates a .pptx presentation from a Markdown outline.

    The Markdown should use H1 (#) for the presentation title,
    H2 (##) for slide titles, and regular text/bullets for slide content.

    Args:
        content_markdown: The Markdown content outline for the presentation.
        output_filename: The name of the output .pptx file.
        theme: The visual theme for the presentation ('default', 'dark', 'minimal').
        task_id: Optional task ID for tracking.

    Returns:
        JSON string with the path to the generated file or an error message.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Theme definitions
        themes = {
            "default": {
                "bg": RGBColor(0xFF, 0xFF, 0xFF),
                "title_color": RGBColor(0x1A, 0x1A, 0x2E),
                "text_color": RGBColor(0x33, 0x33, 0x33),
                "accent": RGBColor(0x16, 0x21, 0x3E),
            },
            "dark": {
                "bg": RGBColor(0x1A, 0x1A, 0x2E),
                "title_color": RGBColor(0xE0, 0xE0, 0xFF),
                "text_color": RGBColor(0xCC, 0xCC, 0xCC),
                "accent": RGBColor(0x53, 0x3F, 0xC8),
            },
            "minimal": {
                "bg": RGBColor(0xFA, 0xFA, 0xFA),
                "title_color": RGBColor(0x22, 0x22, 0x22),
                "text_color": RGBColor(0x44, 0x44, 0x44),
                "accent": RGBColor(0x00, 0x7A, 0xFF),
            }
        }
        t = themes.get(theme, themes["default"])

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        def set_bg(slide, color):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_textbox(slide, text, left, top, width, height,
                        font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            return txBox

        # Parse Markdown into slides
        lines = content_markdown.strip().split('\n')
        slides_data = []
        current_slide = None

        for line in lines:
            if line.startswith('# ') and not line.startswith('## '):
                # Presentation title slide
                slides_data.append({'type': 'title', 'title': line[2:].strip(), 'content': []})
            elif line.startswith('## '):
                if current_slide:
                    slides_data.append(current_slide)
                current_slide = {'type': 'content', 'title': line[3:].strip(), 'content': []}
            elif current_slide and line.strip():
                current_slide['content'].append(line.strip())

        if current_slide:
            slides_data.append(current_slide)

        # Generate slides
        for slide_data in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide, t['bg'])

            if slide_data['type'] == 'title':
                # Title slide
                add_textbox(slide, slide_data['title'],
                            left=1.0, top=2.5, width=11.33, height=1.5,
                            font_size=44, bold=True, color=t['title_color'],
                            align=PP_ALIGN.CENTER)
                # Accent line
                from pptx.util import Pt as PtU
                line_shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Inches(4.5), Inches(4.2), Inches(4.33), Inches(0.05)
                )
                line_shape.fill.solid()
                line_shape.fill.fore_color.rgb = t['accent']
                line_shape.line.fill.background()
            else:
                # Content slide
                # Title bar
                title_bg = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
                )
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = t['accent']
                title_bg.line.fill.background()

                add_textbox(slide, slide_data['title'],
                            left=0.3, top=0.1, width=12.73, height=1.0,
                            font_size=28, bold=True,
                            color=RGBColor(0xFF, 0xFF, 0xFF),
                            align=PP_ALIGN.LEFT)

                # Content
                content_text = '\n'.join(slide_data['content'])
                # Clean up markdown bullets
                content_text = re.sub(r'^[-*•]\s+', '• ', content_text, flags=re.MULTILINE)
                add_textbox(slide, content_text,
                            left=0.5, top=1.4, width=12.33, height=5.8,
                            font_size=20, bold=False, color=t['text_color'])

        # Save
        workspace = os.path.expanduser(os.environ.get("HERMES_WORKSPACE_DIR", "~/.hermes/workspace"))
        os.makedirs(workspace, exist_ok=True)
        output_path = os.path.join(workspace, output_filename)
        prs.save(output_path)

        return json.dumps({
            "success": True,
            "file_path": output_path,
            "slide_count": len(slides_data),
            "message": f"Presentation saved to {output_path} with {len(slides_data)} slides."
        })

    except ImportError:
        return json.dumps({
            "success": False,
            "error": "python-pptx is not installed. Run: pip install python-pptx"
        })
    except Exception as e:
        return json.dumps({
            "success": False,
            "error": str(e)
        })


registry.register(
    name="create_presentation",
    toolset="productivity",
    schema={
        "name": "create_presentation",
        "description": (
            "Generates a .pptx slide deck from a Markdown content outline. "
            "Use H1 (#) for the presentation title slide, H2 (##) for each slide title, "
            "and regular text or bullet points for slide content. "
            "Returns the file path of the generated presentation."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content_markdown": {
                    "type": "string",
                    "description": "The full Markdown content outline for the presentation."
                },
                "output_filename": {
                    "type": "string",
                    "description": "The filename for the output .pptx file (e.g., 'my_presentation.pptx').",
                    "default": "presentation.pptx"
                },
                "theme": {
                    "type": "string",
                    "enum": ["default", "dark", "minimal"],
                    "description": "The visual theme for the presentation.",
                    "default": "default"
                }
            },
            "required": ["content_markdown"]
        }
    },
    handler=lambda args, **kw: create_presentation(
        content_markdown=args.get("content_markdown", ""),
        output_filename=args.get("output_filename", "presentation.pptx"),
        theme=args.get("theme", "default"),
        task_id=kw.get("task_id")
    ),
    check_fn=check_requirements,
    requires_env=[],
)
